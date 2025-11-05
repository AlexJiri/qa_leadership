from pptx import Presentation
from pptx.util import Inches
import datetime
import os

def generate_pptx_report(data):
    prs = Presentation()
    
    # Get data - handle both old and new formats
    title = data.get('title', 'QA Leadership Report')
    date = data.get('date', '')
    kpis = data.get('kpis', [])
    feedback = data.get('feedback', {})
    ai_summary = data.get('ai_summary', '')
    
    # Slide 1 – Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    subtitle_text = f"Date: {date}" if date else f"Generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    subtitle_shape.text = subtitle_text

    # Slide 2 – Key Performance Indicators (KPIs)
    if kpis:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Key Performance Indicators"
        
        # Format KPIs for display - handle both string and object formats
        kpi_list = []
        for kpi in kpis[:10]:
            if isinstance(kpi, str):
                kpi_list.append(f"• {kpi}")
            elif isinstance(kpi, dict):
                # New format: has 'name' and 'how_to_measure'
                if 'name' in kpi and 'how_to_measure' in kpi:
                    kpi_list.append(f"• {kpi.get('name')}: {kpi.get('how_to_measure')}")
                # Old format: has 'kpi' or 'text'
                else:
                    kpi_list.append(f"• {kpi.get('kpi', kpi.get('text', ''))}")
            else:
                kpi_list.append(f"• {str(kpi)}")
        
        kpi_text = "\n".join(kpi_list) if kpi_list else "No KPIs available"
        
        slide.placeholders[1].text = kpi_text[:2000]  # Limit to 2000 chars for slide
    
    # Slide 3 – AI Summary (if available)
    if ai_summary and ai_summary != "—":
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "AI Summary"
        slide.placeholders[1].text = ai_summary[:2000]
    
    # Slide 4 – Feedback Summary (if available)
    if feedback and isinstance(feedback, dict):
        rating = feedback.get('rating', 0)
        comments = feedback.get('comments', [])
        if rating > 0 or comments:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Feedback Summary"
            feedback_text = f"Rating: {rating}/5\n\n"
            if comments:
                feedback_text += "Comments:\n" + "\n".join([f"• {c}" for c in comments[:5]])
            slide.placeholders[1].text = feedback_text

    # Generate filename and save
    timestamp = int(datetime.datetime.now().timestamp())
    filename = f"qa_report_{timestamp}.pptx"
    
    # Ensure generated directory exists
    generated_dir = "static/generated"
    os.makedirs(generated_dir, exist_ok=True)
    
    output_path = os.path.join(generated_dir, filename)
    prs.save(output_path)
    
    return filename
